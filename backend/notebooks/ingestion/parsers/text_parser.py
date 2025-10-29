"""
Text parser for markdown, plain text, and Office documents.
"""

import logging
from typing import Any, Optional

from ..exceptions import ParseError
from .base_parser import BaseParser, ParseResult


class TextParser(BaseParser):
    """Unified text parser for various text-based formats."""

    def __init__(self, logger: Optional[logging.Logger] = None):
        self.logger = logger or logging.getLogger(__name__)

    async def parse(self, file_path: str, metadata: dict[str, Any]) -> ParseResult:
        """
        Parse text file based on extension.

        Args:
            file_path: Path to text file (can be None for in-memory content)
            metadata: File metadata including 'content' for in-memory parsing

        Returns:
            ParseResult with extracted content and metadata
        """
        file_extension = metadata.get("file_extension", "").lower()

        # Route to appropriate parser
        if file_extension == ".md":
            return await self._parse_markdown(file_path, metadata)
        elif file_extension == ".txt":
            return await self._parse_plain_text(file_path, metadata)
        elif file_extension in [".doc", ".docx"]:
            return await self._parse_word_document(file_path, metadata)
        elif file_extension in [".ppt", ".pptx"]:
            return await self._parse_presentation(file_path, metadata)
        else:
            # Generic text handling
            return await self._parse_plain_text(file_path, metadata)

    async def _parse_markdown(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse markdown file directly."""
        self.logger.info(f"Processing markdown file: {file_path}")

        # Check if content is provided in metadata (for in-memory parsing)
        if "content" in metadata and metadata["content"]:
            content = metadata["content"]
        else:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

        text_metadata = {
            "word_count": len(content.split()),
            "char_count": len(content),
            "line_count": len(content.splitlines()),
            "encoding": "utf-8",
            "processing_method": "direct_markdown",
            "is_markdown_file": True,
        }

        # Use original filename for markdown files
        original_filename = metadata.get("filename", "document.md")

        return ParseResult(
            content=content,
            metadata=text_metadata,
            features_available=["content_analysis", "summarization"],
        )

    async def _parse_plain_text(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse plain text file."""
        self.logger.info(f"Processing text file: {file_path}")

        # Check if content is provided in metadata (for in-memory parsing)
        if "content" in metadata and metadata["content"]:
            content = metadata["content"]
            encoding = "utf-8"
        else:
            # Try different encodings
            content = None
            encoding = None
            for enc in ["utf-8", "latin-1", "cp1252"]:
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        content = f.read()
                    encoding = enc
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                raise ParseError("Could not decode text file with any supported encoding")

        text_metadata = {
            "word_count": len(content.split()),
            "char_count": len(content),
            "line_count": len(content.splitlines()),
            "encoding": encoding,
        }

        return ParseResult(
            content=content,
            metadata=text_metadata,
            features_available=["content_analysis", "summarization"],
        )

    async def _parse_word_document(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse Word document (.doc, .docx)."""
        self.logger.info(f"Processing Word document: {file_path}")

        try:
            # Try python-docx for .docx files
            if metadata.get("file_extension", "").lower() == ".docx":
                try:
                    from docx import Document

                    doc = Document(file_path)
                    content = "\n\n".join([para.text for para in doc.paragraphs if para.text])

                    doc_metadata = {
                        "processing_method": "python_docx",
                        "paragraph_count": len(doc.paragraphs),
                        "word_count": len(content.split()),
                    }

                    return ParseResult(
                        content=content,
                        metadata=doc_metadata,
                        features_available=["text_extraction", "document_structure"],
                    )
                except ImportError:
                    self.logger.warning("python-docx not available, falling back to simple message")
                    pass

            # Fallback: return a message that file is uploaded but needs processing
            content = (
                f"Word document '{metadata['filename']}' uploaded successfully. "
                f"Text extraction requires python-docx installation for detailed processing."
            )

            return ParseResult(
                content=content,
                metadata={
                    "file_type": "word_document",
                    "supported_extraction": False,
                },
                features_available=["text_extraction", "formatting_preservation"],
            )

        except Exception as e:
            raise ParseError(f"Word document processing failed: {e}") from e

    async def _parse_presentation(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse PowerPoint presentation (.ppt, .pptx)."""
        self.logger.info(f"Processing presentation: {file_path}")

        try:
            # Try python-pptx for extraction
            try:
                from pptx import Presentation

                prs = Presentation(file_path)

                content = ""
                for i, slide in enumerate(prs.slides):
                    content += f"\n=== Slide {i + 1} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            content += shape.text + "\n"

                ppt_metadata = {
                    "processing_method": "python_pptx",
                    "slide_count": len(prs.slides),
                    "word_count": len(content.split()),
                }

                return ParseResult(
                    content=content,
                    metadata=ppt_metadata,
                    features_available=["text_extraction", "slide_structure"],
                )

            except ImportError:
                self.logger.warning("python-pptx not available, falling back to simple message")
                # Fallback
                content = (
                    f"Presentation '{metadata['filename']}' uploaded successfully. "
                    f"Text extraction requires python-pptx installation for detailed processing."
                )

                return ParseResult(
                    content=content,
                    metadata={
                        "file_type": "presentation",
                        "supported_extraction": False,
                    },
                    features_available=[
                        "slide_extraction",
                        "text_extraction",
                        "image_extraction",
                    ],
                )

        except Exception as e:
            raise ParseError(f"Presentation processing failed: {e}") from e
