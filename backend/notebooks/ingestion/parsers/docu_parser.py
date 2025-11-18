"""
Document parser using MinerU API for PDFs, Word, and PowerPoint files, with PyMuPDF fallback for PDFs.
"""

import base64
import json
import logging
import os
import subprocess
import tempfile
import time
from typing import Any, Optional

import requests

from ...utils.helpers import clean_title
from ..exceptions import ParseError
from .base_parser import BaseParser, ParseResult

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None


class DocuParser(BaseParser):
    """Document parser with MinerU for PDFs, Word, PowerPoint and PyMuPDF fallback for PDFs."""

    def __init__(
        self,
        mineru_base_url: str,
        logger: Optional[logging.Logger] = None,
    ):
        # Normalize URL
        if not str(mineru_base_url).lower().startswith(("http://", "https://")):
            mineru_base_url = f"http://{mineru_base_url}"
        self.mineru_base_url = mineru_base_url.rstrip("/")
        self.mineru_parse_endpoint = f"{self.mineru_base_url}/file_parse"
        self.logger = logger or logging.getLogger(__name__)

        # No CSS needed; using LibreOffice for conversions

    def _convert_pptx_to_pdf(self, filepath: str) -> str:
        """Convert PPT/PPTX to PDF using LibreOffice headless and return temp PDF path."""
        outdir = tempfile.mkdtemp(prefix="lo_ppt_", suffix="_to_pdf")
        try:
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                outdir,
                filepath,
            ]
            subprocess.run(cmd, check=True)
            out_pdf = os.path.join(
                outdir, os.path.splitext(os.path.basename(filepath))[0] + ".pdf"
            )
            if not os.path.exists(out_pdf):
                raise ParseError(
                    f"PPT/PPTX to PDF conversion produced no output for {filepath}"
                )
            return out_pdf
        except Exception as e:
            raise ParseError(f"PPT/PPTX to PDF conversion failed: {e}") from e

    def _convert_docx_to_pdf(self, filepath: str) -> str:
        """Convert DOC/DOCX to PDF using LibreOffice headless and return temp PDF path."""
        outdir = tempfile.mkdtemp(prefix="lo_doc_", suffix="_to_pdf")
        try:
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                outdir,
                filepath,
            ]
            subprocess.run(cmd, check=True)
            out_pdf = os.path.join(
                outdir, os.path.splitext(os.path.basename(filepath))[0] + ".pdf"
            )
            if not os.path.exists(out_pdf):
                raise ParseError(
                    f"DOC/DOCX to PDF conversion produced no output for {filepath}"
                )
            return out_pdf
        except Exception as e:
            raise ParseError(f"DOC/DOCX to PDF conversion failed: {e}") from e

    # Removed HTML-based PPTX/DOCX helpers; LibreOffice handles formatting directly

    async def parse(self, file_path: str, metadata: dict[str, Any]) -> ParseResult:
        """
        Parse document file (PDF, Word, PowerPoint) using MinerU, with fallbacks.

        Args:
            file_path: Path to document file
            metadata: File metadata (filename, extension, etc.)

        Returns:
            ParseResult with content, metadata, and optional mineru_extraction_result
        """
        file_extension = metadata.get("file_extension", "").lower()
        is_pdf = file_extension == ".pdf"
        is_word = file_extension in [".docx", ".doc"]
        is_powerpoint = file_extension in [".pptx", ".ppt"]

        # Try MinerU first (no health check - assume it's available)
        try:
            return await self._parse_with_mineru(file_path, metadata)
        except Exception as e:
            self.logger.warning(
                f"MinerU parsing failed: {e}. Attempting fallback parser."
            )

            # Fall back to appropriate parser based on file type
            if is_pdf:
                self.logger.info("Falling back to PyMuPDF for PDF")
                return await self._parse_with_pymupdf(file_path, metadata)
            elif is_word:
                self.logger.info("Falling back to python-docx for Word document")
                return await self._parse_with_docx(file_path, metadata)
            elif is_powerpoint:
                self.logger.info("Falling back to python-pptx for PowerPoint")
                return await self._parse_with_pptx(file_path, metadata)
            else:
                raise ParseError(
                    f"MinerU parsing failed for {file_extension} file and no fallback parser is available. Error: {e}"
                )

    def _check_mineru_health(self) -> bool:
        """Check if MinerU API is available."""
        try:
            response = requests.get(f"{self.mineru_base_url}/docs", timeout=10)
            return response.status_code == 200
        except Exception as e:
            self.logger.warning(f"MinerU health check failed: {e}")
            return False

    async def _parse_with_mineru(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse document (PDF, Word, PowerPoint) using MinerU API."""
        file_extension = metadata.get("file_extension", "").lower()

        # Determine file type for logging
        file_type_map = {
            ".pdf": "PDF",
            ".docx": "Word",
            ".doc": "Word",
            ".pptx": "PowerPoint",
            ".ppt": "PowerPoint",
        }
        file_type = file_type_map.get(file_extension, "Document")
        self.logger.info(f"Starting MinerU parsing of {file_type} file: {file_path}")
        start_time = time.time()

        # Convert PPTX/DOCX to PDF if needed
        temp_pdf_path = None
        actual_file_path = file_path

        try:
            if file_extension in [".pptx", ".ppt"]:
                self.logger.info(f"Converting {file_type} to PDF for MinerU processing")
                temp_pdf_path = self._convert_pptx_to_pdf(file_path)
                actual_file_path = temp_pdf_path
            elif file_extension in [".docx", ".doc"]:
                self.logger.info(f"Converting {file_type} to PDF for MinerU processing")
                temp_pdf_path = self._convert_docx_to_pdf(file_path)
                actual_file_path = temp_pdf_path

            # Generate clean filename
            default_filename = (
                f"document{file_extension}" if file_extension else "document.pdf"
            )
            original_filename = metadata.get("filename", default_filename)
            base_title = (
                original_filename.rsplit(".", 1)[0]
                if "." in original_filename
                else original_filename
            )
            clean_pdf_title = clean_title(base_title)

            # Call MinerU API with the actual file (original PDF or converted PDF)
            mineru_result = self._call_mineru_api(actual_file_path)

            # Extract results
            results = mineru_result.get("results", {})
            if not results:
                raise ParseError("No results returned from MinerU API")

            # Get first document result
            doc_key = next(iter(results.keys()))
            doc_result = results[doc_key]

            # Extract markdown content and images
            md_content = doc_result.get("md_content", "")
            images = doc_result.get("images", {})

            # Validate MinerU output
            if not md_content:
                self.logger.warning(
                    f"MinerU returned no markdown content for {original_filename}"
                )

            if not images:
                self.logger.info(
                    f"MinerU extracted no images from {original_filename} "
                    "(document may not contain images or images may not be extractable)"
                )
            else:
                self.logger.info(
                    f"MinerU extracted {len(images)} images from {original_filename}: "
                    f"{list(images.keys())}"
                )

            # Save to temporary directory
            temp_dir = tempfile.mkdtemp(suffix="_mineru_output")

            # Save markdown content
            md_file_path = os.path.join(temp_dir, f"{doc_key}.md")
            with open(md_file_path, "w", encoding="utf-8") as f:
                f.write(md_content)

            # Save images
            image_files = []
            failed_images = []
            for img_name, img_data in images.items():
                try:
                    if img_data.startswith("data:image/"):
                        # Decode base64 images
                        header, data = img_data.split(",", 1)
                        img_bytes = base64.b64decode(data)

                        img_path = os.path.join(temp_dir, img_name)
                        with open(img_path, "wb") as f:
                            f.write(img_bytes)
                        image_files.append(img_path)
                        self.logger.debug(
                            f"Saved image {img_name}: {len(img_bytes)} bytes"
                        )
                    else:
                        self.logger.warning(
                            f"Skipping image {img_name}: invalid data format "
                            f"(expected data:image/*, got: {img_data[:50]}...)"
                        )
                        failed_images.append(img_name)
                except Exception as e:
                    self.logger.error(
                        f"Failed to decode/save image {img_name}: {str(e)}"
                    )
                    failed_images.append(img_name)

            # Log summary
            if failed_images:
                self.logger.warning(
                    f"Failed to save {len(failed_images)} images: {failed_images}"
                )

            self.logger.info(
                f"Created {len(os.listdir(temp_dir))} files in temp directory: "
                f"1 markdown file, {len(image_files)} images"
            )

            # Get document metadata
            doc_metadata = self._extract_document_metadata(
                file_path, metadata, mineru_result
            )
            doc_metadata["has_mineru_extraction"] = True
            doc_metadata["has_markdown_content"] = bool(md_content)
            doc_metadata["file_type"] = file_type.lower()

            # Calculate features
            features_available = ["advanced_pdf_extraction", "markdown_conversion"]
            if images:
                features_available.append("image_extraction")
            if "table" in md_content.lower() or "|" in md_content:
                features_available.append("table_extraction")
            if any(marker in md_content for marker in ["$$", "$", "\\("]):
                features_available.append("formula_extraction")
            features_available.append("layout_analysis")

            duration = time.time() - start_time
            self.logger.info(f"MinerU parsing completed in {duration:.2f} seconds")

            # Return result with MinerU extraction info
            return ParseResult(
                content="",  # Empty content since MinerU files contain everything
                metadata=doc_metadata,
                features_available=features_available,
                mineru_extraction_result={
                    "success": True,
                    "temp_mineru_dir": temp_dir,
                    "clean_title": clean_pdf_title,
                },
            )

        finally:
            # Clean up temporary PDF file if it was created
            if temp_pdf_path and os.path.exists(temp_pdf_path):
                try:
                    os.remove(temp_pdf_path)
                    self.logger.debug(f"Cleaned up temporary PDF: {temp_pdf_path}")
                except Exception as e:
                    self.logger.warning(f"Failed to clean up temporary PDF: {e}")

    async def _parse_with_pymupdf(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse PDF using PyMuPDF fallback."""
        self.logger.info(f"Using PyMuPDF fallback for {file_path}")

        if not fitz:
            raise ParseError("PyMuPDF (fitz) is not available")

        doc = fitz.open(file_path)
        content = ""

        # Extract metadata
        pdf_metadata = {
            "page_count": doc.page_count,
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "creation_date": doc.metadata.get("creationDate", ""),
            "modification_date": doc.metadata.get("modDate", ""),
            "processing_method": "pymupdf_fallback",
        }

        # Extract text from all pages
        for page_num in range(doc.page_count):
            page = doc[page_num]
            content += f"\n=== Page {page_num + 1} ===\n"
            page_text = page.get_text()
            content += page_text

        doc.close()

        # Check if extraction was successful
        if not content.strip():
            content = (
                f"PDF document '{metadata['filename']}' appears to be image-based or empty. "
                f"Text extraction may require OCR processing."
            )

        return ParseResult(
            content=content,
            metadata=pdf_metadata,
            features_available=[
                "advanced_pdf_extraction",
                "figure_extraction",
                "table_extraction",
            ],
        )

    async def _parse_with_docx(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse Word document using python-docx fallback."""
        self.logger.info(f"Using python-docx fallback for {file_path}")

        try:
            # Try python-docx for .docx files
            if metadata.get("file_extension", "").lower() == ".docx":
                try:
                    from docx import Document

                    doc = Document(file_path)
                    content = "\n\n".join(
                        [para.text for para in doc.paragraphs if para.text]
                    )

                    doc_metadata = {
                        "processing_method": "python_docx_fallback",
                        "paragraph_count": len(doc.paragraphs),
                        "word_count": len(content.split()),
                    }

                    return ParseResult(
                        content=content,
                        metadata=doc_metadata,
                        features_available=["text_extraction", "document_structure"],
                    )
                except ImportError:
                    self.logger.warning(
                        "python-docx not available, returning placeholder"
                    )

            # Fallback: return a message that file is uploaded but needs processing
            content = (
                f"Word document '{metadata['filename']}' uploaded successfully. "
                f"Text extraction requires python-docx installation for detailed processing."
            )

            return ParseResult(
                content=content,
                metadata={
                    "file_type": "word_document",
                    "supported_extraction": False,
                    "processing_method": "fallback_placeholder",
                },
                features_available=["text_extraction", "formatting_preservation"],
            )

        except Exception as e:
            raise ParseError(f"Word document fallback processing failed: {e}") from e

    async def _parse_with_pptx(
        self, file_path: str, metadata: dict[str, Any]
    ) -> ParseResult:
        """Parse PowerPoint presentation using python-pptx fallback."""
        self.logger.info(f"Using python-pptx fallback for {file_path}")

        try:
            # Try python-pptx for extraction
            try:
                from pptx import Presentation

                prs = Presentation(file_path)

                content = ""
                for i, slide in enumerate(prs.slides):
                    content += f"\n=== Slide {i + 1} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            content += shape.text + "\n"

                ppt_metadata = {
                    "processing_method": "python_pptx_fallback",
                    "slide_count": len(prs.slides),
                    "word_count": len(content.split()),
                }

                return ParseResult(
                    content=content,
                    metadata=ppt_metadata,
                    features_available=["text_extraction", "slide_structure"],
                )

            except ImportError:
                self.logger.warning("python-pptx not available, returning placeholder")
                # Fallback
                content = (
                    f"Presentation '{metadata['filename']}' uploaded successfully. "
                    f"Text extraction requires python-pptx installation for detailed processing."
                )

                return ParseResult(
                    content=content,
                    metadata={
                        "file_type": "presentation",
                        "supported_extraction": False,
                        "processing_method": "fallback_placeholder",
                    },
                    features_available=[
                        "slide_extraction",
                        "text_extraction",
                        "image_extraction",
                    ],
                )

        except Exception as e:
            raise ParseError(f"Presentation fallback processing failed: {e}") from e

    def _call_mineru_api(self, file_path: str) -> dict[str, Any]:
        """Call MinerU API to parse document (PDF, Word, PowerPoint)."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document file not found: {file_path}")

        try:
            # Detect file type and set appropriate MIME type
            file_extension = os.path.splitext(file_path)[1].lower()
            mime_type_map = {
                ".pdf": "application/pdf",
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".doc": "application/msword",
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".ppt": "application/vnd.ms-powerpoint",
            }
            mime_type = mime_type_map.get(file_extension, "application/octet-stream")

            # Prepare multipart form data
            files = {
                "files": (
                    os.path.basename(file_path),
                    open(file_path, "rb"),
                    mime_type,
                )
            }

            data = {
                "output_dir": "./output",
                "lang_list": ["ch"],
                "backend": "pipeline",
                "parse_method": "auto",
                "formula_enable": True,
                "table_enable": True,
                "server_url": "10.218.163.144",
                "return_md": True,
                "return_middle_json": False,
                "return_model_output": False,
                "return_content_list": False,
                "return_images": True,
                "response_format_zip": False,
                "start_page_id": 0,
                "end_page_id": 99999,
            }

            # Make API request
            response = requests.post(
                self.mineru_parse_endpoint,
                files=files,
                data=data,
                timeout=300,  # 5 minute timeout
            )

            # Close file handle
            files["files"][1].close()

            response.raise_for_status()
            return response.json()

        except requests.exceptions.RequestException as e:
            self.logger.error(f"MinerU API request failed: {e}")
            raise ParseError(f"MinerU API request failed: {e}") from e
        except json.JSONDecodeError as e:
            self.logger.error(f"Failed to parse MinerU API response: {e}")
            raise ParseError(f"Invalid JSON response from MinerU API: {e}") from e

    def _extract_document_metadata(
        self, file_path: str, metadata: dict[str, Any], mineru_result: dict[str, Any]
    ) -> dict[str, Any]:
        """Extract document metadata (PDF, Word, PowerPoint)."""
        file_extension = metadata.get("file_extension", "").lower()
        is_pdf = file_extension == ".pdf"

        base_metadata = {
            "processing_method": "mineru_api",
            "api_version": mineru_result.get("version", "unknown"),
            "backend": mineru_result.get("backend", "pipeline"),
        }

        # For PDF files, try to extract additional metadata using PyMuPDF
        if is_pdf:
            try:
                if fitz:
                    doc = fitz.open(file_path)
                    base_metadata.update(
                        {
                            "page_count": doc.page_count,
                            "title": doc.metadata.get("title", ""),
                            "author": doc.metadata.get("author", ""),
                            "creation_date": doc.metadata.get("creationDate", ""),
                            "modification_date": doc.metadata.get("modDate", ""),
                        }
                    )
                    doc.close()
            except Exception as e:
                self.logger.warning(f"Could not extract PDF metadata: {e}")
                base_metadata["metadata_error"] = str(e)

        return base_metadata
