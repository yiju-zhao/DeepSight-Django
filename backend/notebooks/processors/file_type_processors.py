"""
File Type Processors - Handle specific file type processing logic.
Extracted from upload_processor.py for better maintainability.
"""
import os
import tempfile
import subprocess
import logging
import time
import requests
import json
import base64
from typing import Dict, Any
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from ..utils.helpers import clean_title
except ImportError:
    clean_title = None


class FileTypeProcessors:
    """Handle file type specific processing logic."""
    
    def __init__(self, mineru_base_url: str, xinference_url: str = None, model_uid: str = None, logger=None):
        # Normalize: ensure URL has a scheme
        if not str(mineru_base_url).lower().startswith(("http://", "https://")):
            mineru_base_url = f"http://{mineru_base_url}"
        self.mineru_base_url = mineru_base_url.rstrip('/')
        self.mineru_parse_endpoint = f"{self.mineru_base_url}/file_parse"
        self.xinference_url = xinference_url or os.getenv('XINFERENCE_URL', 'http://localhost:9997')
        self.model_uid = model_uid or os.getenv('XINFERENCE_WHISPER_MODEL_UID', 'whisper-large-v3-turbo')
        self.logger = logger or logging.getLogger(__name__)
    
    def log_operation(self, operation: str, details: str = "", level: str = "info"):
        """Log operations with consistent formatting."""
        message = f"[file_type_processors] {operation}"
        if details:
            message += f": {details}"
        getattr(self.logger, level)(message)

    def check_mineru_health(self) -> bool:
        """Check if MinerU API is available."""
        try:
            response = requests.get(f"{self.mineru_base_url}/docs", timeout=10)
            return response.status_code == 200
        except Exception as e:
            self.log_operation("mineru_health_check_failed", f"MinerU API health check failed: {e}", "warning")
            return False

    async def process_file_by_type(self, file_path: str, file_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process file based on its type."""
        file_extension = file_metadata.get('file_extension', '').lower()
        
        if file_extension == '.pdf':
            return self.process_pdf_mineru(file_path, file_metadata)
        elif file_extension in ['.mp3', '.wav', '.m4a']:
            return await self.process_audio_immediate(file_path, file_metadata)
        elif file_extension in [".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".3gp", ".ogv", ".m4v"]:
            return await self.process_video_immediate(file_path, file_metadata)
        elif file_extension == ".md":
            return self.process_markdown_direct(file_path, file_metadata)
        elif file_extension == ".txt":
            return self.process_text_immediate(file_path, file_metadata)
        elif file_extension in [".ppt", ".pptx", ".doc", ".docx"]:
            return self.process_presentation_immediate(file_path, file_metadata)
        else:
            return {
                'content': f"File type {file_extension} is supported but no immediate processing available.",
                'metadata': {},
                'features_available': [],
                'processing_time': 'immediate'
            }

    def process_pdf_mineru(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """PDF text extraction using MinerU API."""
        try:
            self.log_operation("pdf_mineru_start", f"Starting MinerU processing of {file_path}")
            start_time = time.time()

            # Check if MinerU API is available
            if not self.check_mineru_health():
                self.log_operation("pdf_mineru_unavailable", "MinerU API not available, falling back to PyMuPDF")
                return self.process_pdf_pymupdf_fallback(file_path, file_metadata)

            # Generate clean filename for PDF
            original_filename = file_metadata.get('filename', 'document')
            base_title = original_filename.rsplit('.', 1)[0] if '.' in original_filename else original_filename
            clean_pdf_title = clean_title(base_title) if clean_title else base_title

            # Call MinerU API
            mineru_result = self._call_mineru_api(file_path)
            
            # Extract the results for the first (and likely only) document
            results = mineru_result.get('results', {})
            if not results:
                raise Exception("No results returned from MinerU API")
            
            # Get the first document result (assuming single PDF)
            doc_key = next(iter(results.keys()))
            doc_result = results[doc_key]
            
            # Extract markdown content
            md_content = doc_result.get('md_content', '')
            
            # Extract images if available
            images = doc_result.get('images', {})
            
            # Save content to temporary file for processing
            temp_dir = tempfile.mkdtemp(suffix='_mineru_output')
            
            # Save markdown content
            md_file_path = os.path.join(temp_dir, f"{doc_key}.md")
            with open(md_file_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            # Save images if any
            image_files = []
            for img_name, img_data in images.items():
                if img_data.startswith('data:image/'):
                    # Handle base64 encoded images
                    header, data = img_data.split(',', 1)
                    img_bytes = base64.b64decode(data)
                    
                    img_path = os.path.join(temp_dir, img_name)
                    with open(img_path, 'wb') as f:
                        f.write(img_bytes)
                    image_files.append(img_path)
            
            self.log_operation("pdf_mineru_files", f"Created {len(os.listdir(temp_dir))} files: {[os.path.basename(f) for f in [md_file_path] + image_files]}")

            # Get basic PDF metadata using PyMuPDF if available
            try:
                if fitz:
                    doc = fitz.open(file_path)
                    pdf_metadata = {
                        'page_count': doc.page_count,
                        'title': doc.metadata.get('title', ''),
                        'author': doc.metadata.get('author', ''),
                        'creation_date': doc.metadata.get('creationDate', ''),
                        'modification_date': doc.metadata.get('modDate', ''),
                        'processing_method': 'mineru_api',
                        'api_version': mineru_result.get('version', 'unknown'),
                        'backend': mineru_result.get('backend', 'pipeline'),
                        'has_markdown_content': bool(md_content),
                        'extracted_images_count': len(images),
                        'has_mineru_extraction': temp_dir is not None
                    }
                    doc.close()
                else:
                    pdf_metadata = {
                        'processing_method': 'mineru_api',
                        'api_version': mineru_result.get('version', 'unknown'),
                        'backend': mineru_result.get('backend', 'pipeline'),
                        'has_markdown_content': bool(md_content),
                        'extracted_images_count': len(images),
                        'has_mineru_extraction': temp_dir is not None
                    }
            except Exception as e:
                self.log_operation("pdf_metadata_warning", f"Could not extract PDF metadata: {e}", "warning")
                pdf_metadata = {
                    'processing_method': 'mineru_api',
                    'api_version': mineru_result.get('version', 'unknown'),
                    'metadata_error': str(e),
                    'has_markdown_content': bool(md_content),
                    'extracted_images_count': len(images),
                    'has_mineru_extraction': temp_dir is not None
                }

            # For MinerU PDFs, we don't store content separately since MinerU files contain everything
            summary_content = ""

            end_time = time.time()
            duration = end_time - start_time
            self.log_operation("pdf_mineru_completed", f"MinerU processing completed in {duration:.2f} seconds")

            # Calculate features available
            features_available = ['advanced_pdf_extraction', 'markdown_conversion']
            if images:
                features_available.append('image_extraction')
            if 'table' in md_content.lower() or '|' in md_content:
                features_available.append('table_extraction')
            if any(marker in md_content for marker in ['$$', '$', '\\\\(']):
                features_available.append('formula_extraction')
            features_available.append('layout_analysis')

            result = {
                'content': summary_content,
                'content_filename': f"{clean_pdf_title}.md", 
                'metadata': pdf_metadata,
                'features_available': features_available,
                'processing_time': f'{duration:.2f}s',
                'skip_content_file': True  # Flag to skip creating extracted_content.md since MinerU provides better content
            }

            # Add mineru extraction result for post-processing
            if temp_dir:
                result['marker_extraction_result'] = {
                    'success': True,
                    'temp_marker_dir': temp_dir,
                    'clean_title': clean_pdf_title
                }

            return result

        except Exception as e:
            self.log_operation("pdf_mineru_error", f"MinerU processing failed: {e}", "warning")
            # Fallback to PyMuPDF if MinerU fails
            return self.process_pdf_pymupdf_fallback(file_path, file_metadata)

    def process_pdf_pymupdf_fallback(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Fallback PDF text extraction using PyMuPDF when MinerU is not available."""
        try:
            self.log_operation("pdf_pymupdf_fallback", f"Using PyMuPDF fallback for {file_path}")

            if not fitz:
                raise Exception("PyMuPDF (fitz) is not available")

            doc = fitz.open(file_path)
            content = ""

            # Extract basic metadata
            pdf_metadata = {
                'page_count': doc.page_count,
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', ''),
                'processing_method': 'pymupdf_fallback'
            }

            # Extract text from all pages
            for page_num in range(doc.page_count):
                page = doc[page_num]
                content += f"\\n=== Page {page_num + 1} ===\\n"
                page_text = page.get_text()
                content += page_text

            doc.close()

            # Check if content extraction was successful
            if not content.strip():
                content = f"PDF document '{file_metadata['filename']}' appears to be image-based or empty. Text extraction may require OCR processing."

            return {
                "content": content,
                "metadata": pdf_metadata,
                "features_available": [
                    "advanced_pdf_extraction",
                    "figure_extraction",
                    "table_extraction",
                ],
                "processing_time": "immediate",
            }

        except Exception as e:
            raise Exception(f"PDF processing failed: {str(e)}")

    def _call_mineru_api(
        self,
        file_path: str,
        output_dir: str = "./output",
        lang_list: list = None,
        backend: str = "pipeline",
        parse_method: str = "auto",
        formula_enable: bool = True,
        table_enable: bool = True,
        server_url: str = "10.218.163.144",
        return_md: bool = True,
        return_middle_json: bool = False,
        return_model_output: bool = False,
        return_content_list: bool = False,
        return_images: bool = True,
        response_format_zip: bool = False,
        start_page_id: int = 0,
        end_page_id: int = 99999
    ) -> Dict[str, Any]:
        """Call the MinerU API to parse a PDF file."""
        if lang_list is None:
            lang_list = ['ch']
            
        # Validate file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")
        
        try:
            # Prepare the multipart form data
            files = {
                'files': (os.path.basename(file_path), open(file_path, 'rb'), 'application/pdf')
            }
            
            data = {
                'output_dir': output_dir,
                'lang_list': lang_list,
                'backend': backend,
                'parse_method': parse_method,
                'formula_enable': formula_enable,
                'table_enable': table_enable,
                'server_url': server_url,
                'return_md': return_md,
                'return_middle_json': return_middle_json,
                'return_model_output': return_model_output,
                'return_content_list': return_content_list,
                'return_images': return_images,
                'response_format_zip': response_format_zip,
                'start_page_id': start_page_id,
                'end_page_id': end_page_id
            }
            
            # Make the API request
            response = requests.post(
                self.mineru_parse_endpoint,
                files=files,
                data=data,
                timeout=300  # 5 minute timeout for large PDFs
            )
            
            # Close the file handle
            files['files'][1].close()
            
            response.raise_for_status()
            
            # Parse the JSON response
            result = response.json()
            
            return result
            
        except requests.exceptions.RequestException as e:
            self.log_operation("mineru_api_request_failed", f"MinerU API request failed: {e}", "error")
            raise Exception(f"MinerU API request failed: {e}")
        except json.JSONDecodeError as e:
            self.log_operation("mineru_api_response_error", f"Failed to parse MinerU API response: {e}", "error")
            raise Exception(f"Invalid JSON response from MinerU API: {e}")
        except Exception as e:
            self.log_operation("mineru_api_call_error", f"MinerU API call failed: {e}", "error")
            raise Exception(f"MinerU API call failed: {e}")

    async def process_audio_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Quick audio transcription using Xinference."""
        try:
            # Always attempt transcription with Xinference - no pre-check needed

            # Transcribe audio using Xinference
            from ..services.transcription_service import TranscriptionService
            transcription_service = TranscriptionService(self.xinference_url, self.model_uid, clean_title, self.logger)
            transcript_content, transcript_filename = await transcription_service.transcribe_audio_video(
                file_path, file_metadata['filename']
            )

            # Get basic audio info
            audio_metadata = self._get_audio_metadata(file_path)
            audio_metadata.update({
                'transcript_filename': transcript_filename,
                'has_transcript': True,
            })

            return {
                "content": transcript_content,
                "metadata": audio_metadata,
                "features_available": [
                    "speaker_diarization",
                    "sentiment_analysis",
                    "advanced_audio_analysis",
                ],
                "processing_time": "immediate",
                "transcript_filename": transcript_filename
            }

        except Exception as e:
            raise Exception(f"Audio processing failed: {str(e)}")

    async def process_video_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Enhanced video processing with optional transcription."""
        try:
            # Extract audio from video for transcription
            audio_path = tempfile.mktemp(suffix='.wav')

            cmd = [
                'ffmpeg', '-i', file_path, '-vn', '-acodec', 'pcm_s16le',
                '-ar', '16000', '-ac', '1', '-y', audio_path
            ]

            result = subprocess.run(cmd, capture_output=True, text=True)

            # Initialize content parts
            content_parts = []

            # Process audio transcription if available
            transcript_filename = None
            has_transcript = False

            # Always generate a transcript filename based on the video name
            base_title = Path(file_metadata['filename']).stem
            cleaned_title = clean_title(base_title) if clean_title else base_title
            transcript_filename = f"{cleaned_title}.md"

            if result.returncode == 0 and os.path.exists(audio_path):
                try:
                    from ..services.transcription_service import TranscriptionService
                    transcription_service = TranscriptionService(self.xinference_url, self.model_uid, clean_title, self.logger)
                    transcript_content, _ = await transcription_service.transcribe_audio_video(
                        audio_path, file_metadata['filename']
                    )
                    content_parts.append(f"# Transcription\\n\\n{transcript_content}")
                    has_transcript = True
                except Exception as e:
                    self.log_operation("video_transcription_error", f"Transcription failed: {e}", "warning")
                finally:
                    # Clean up extracted audio
                    if os.path.exists(audio_path):
                        os.unlink(audio_path)
            else:
                # If no audio or transcription failed
                if result.returncode != 0:
                    content_parts.append(f"# Video: {file_metadata['filename']}\\n\\nNo audio track found or audio extraction failed.")
                else:
                    content_parts.append(f"# Video: {file_metadata['filename']}\\n\\nAudio transcription service not available.")

            # Get video metadata and add transcript info
            video_metadata = self._get_video_metadata(file_path)
            video_metadata.update({
                'transcript_filename': transcript_filename,
                'has_transcript': has_transcript,
                'has_audio': result.returncode == 0,
            })

            # Combine content
            final_content = "\\n\\n".join(content_parts)

            return {
                'content': final_content,
                'metadata': video_metadata,
                'features_available': ['frame_extraction', 'scene_analysis', 'speaker_diarization', 'video_analysis'],
                'processing_time': 'immediate',
                'transcript_filename': transcript_filename if has_transcript else None
            }

        except Exception as e:
            raise Exception(f"Video processing failed: {str(e)}")

    def process_markdown_direct(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process markdown files directly, returning the original file without additional processing."""
        try:
            self.log_operation("markdown_direct_processing", f"Processing markdown file directly: {file_path}")
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            # For markdown files, we don't use the marker extraction logic.
            # We just return the content and use the original filename.
            # The content will be indexed directly without creating additional files.
            text_metadata = {
                "word_count": len(content.split()),
                "char_count": len(content),
                "line_count": len(content.splitlines()),
                "encoding": "utf-8",
                "processing_method": "direct_markdown",
                "is_markdown_file": True,
            }

            # Use the original filename directly - no processing needed for .md files
            original_filename = file_metadata['filename']

            return {
                "content": content,
                "metadata": text_metadata,
                "features_available": ["content_analysis", "summarization"],
                "processing_time": "immediate",
                "content_filename": original_filename,  # Use original filename
                "use_original_file": True   # Flag to indicate we should use the original uploaded file
            }
        except Exception as e:
            self.log_operation("markdown_direct_error", f"Error processing markdown file directly: {e}", "error")
            raise Exception(f"Markdown processing failed: {str(e)}")

    def process_text_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process text files immediately."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            text_metadata = {
                "word_count": len(content.split()),
                "char_count": len(content),
                "line_count": len(content.splitlines()),
                "encoding": "utf-8",
            }

            # Check if this is a pasted text file
            is_pasted_text = file_metadata.get('source_type') == 'pasted_text' or file_metadata.get('original_filename') == 'pasted_text.md'
            
            result = {
                "content": content,
                "metadata": text_metadata,
                "features_available": ["content_analysis", "summarization"],
                "processing_time": "immediate",
            }

            # For pasted text, use specific filename to store in content folder
            if is_pasted_text:
                result["content_filename"] = "pasted_text.md"

            return result

        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ["latin-1", "cp1252"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()

                    text_metadata = {
                        "word_count": len(content.split()),
                        "char_count": len(content),
                        "line_count": len(content.splitlines()),
                        "encoding": encoding,
                    }

                    return {
                        "content": content,
                        "metadata": text_metadata,
                        "features_available": ["content_analysis", "summarization"],
                        "processing_time": "immediate",
                    }
                except UnicodeDecodeError:
                    continue

            raise Exception(f"Could not decode text file with any supported encoding")
        except Exception as e:
            raise Exception(f"Text processing failed: {str(e)}")

    def process_presentation_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process presentation files with fallback message."""
        try:
            # Try to extract text using python-pptx for simple extraction
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                
                content = ""
                for i, slide in enumerate(prs.slides):
                    content += f"\\n=== Slide {i + 1} ===\\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\\n"
                
                return {
                    'content': content,
                    'metadata': {
                        'processing_method': 'python_pptx',
                        'slide_count': len(prs.slides)
                    },
                    'features_available': ['text_extraction', 'slide_structure'],
                    'processing_time': 'immediate'
                }
            
            except ImportError:
                # Fallback if python-pptx not available
                ppt_metadata = {"file_type": "presentation", "supported_extraction": False}
                content = f"Presentation file '{file_metadata['filename']}' uploaded successfully. Text extraction requires python-pptx installation."

                return {
                    "content": content,
                    "metadata": ppt_metadata,
                    "features_available": [
                        "slide_extraction",
                        "text_extraction", 
                        "image_extraction",
                    ],
                    "processing_time": "immediate",
                }
        
        except Exception as e:
            raise Exception(f"Presentation processing failed: {str(e)}")

    def _get_audio_metadata(self, file_path: str) -> Dict:
        """Extract audio metadata using ffprobe."""
        try:
            cmd = [
                "ffprobe",
                "-v",
                "quiet",
                "-print_format",
                "json",
                "-show_format",
                file_path,
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, check=True)

            import json

            data = json.loads(result.stdout)
            format_info = data.get("format", {})

            return {
                "duration": float(format_info.get("duration", 0)),
                "bitrate": int(format_info.get("bit_rate", 0)),
                "size": int(format_info.get("size", 0)),
                "format_name": format_info.get("format_name", "unknown"),
            }
        except Exception:
            return {"error": "Could not extract audio metadata"}

    def _get_video_metadata(self, file_path: str) -> Dict:
        """Extract video metadata using ffprobe."""
        try:
            cmd = [
                "ffprobe",
                "-v",
                "quiet",
                "-print_format",
                "json",
                "-show_streams",
                "-show_format",
                file_path,
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, check=True)

            import json

            data = json.loads(result.stdout)

            # Get video stream info
            video_stream = next(
                (s for s in data.get("streams", []) if s.get("codec_type") == "video"),
                {},
            )
            format_info = data.get("format", {})

            return {
                "duration": float(format_info.get("duration", 0)),
                "resolution": f"{video_stream.get('width', 0)}x{video_stream.get('height', 0)}",
                "fps": video_stream.get("r_frame_rate", "0/0"),
                "codec": video_stream.get("codec_name", "unknown"),
                "format_name": format_info.get("format_name", "unknown"),
                "size": int(format_info.get("size", 0)),
            }
        except Exception:
            return {'error': 'Could not extract video metadata'}
