"""
File Processors - Handle file type specific processing logic
"""
import os
import tempfile
import time
import logging
import requests
import json
import base64
from typing import Dict, Any, Optional
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

logger = logging.getLogger(__name__)


class FileProcessor:
    """Handle file type specific processing"""
    
    def __init__(self, mineru_base_url: str = "http://10.218.163.144:8008"):
        self.logger = logging.getLogger(f"{__name__}.file_processor")
        
        # Initialize whisper model lazily
        self._whisper_model = None
        # MinerU API configuration
        self.mineru_base_url = mineru_base_url.rstrip('/')
        self.mineru_parse_endpoint = f"{self.mineru_base_url}/file_parse"
    
    @property 
    def whisper_model(self):
        """Lazy load whisper model."""
        if self._whisper_model is None:
            try:
                import faster_whisper
                device = self._detect_device()
                self._whisper_model = faster_whisper.WhisperModel("large-v3-turbo", device=device)
                self.logger.info(f"Loaded Whisper model on {device}")
            except ImportError:
                self.logger.warning("faster-whisper not available")
                self._whisper_model = False
        return self._whisper_model

    def check_mineru_health(self) -> bool:
        """Check if MinerU API is available."""
        try:
            response = requests.get(f"{self.mineru_base_url}/docs", timeout=10)
            return response.status_code == 200
        except Exception as e:
            self.logger.warning(f"MinerU API health check failed: {e}")
            return False

    def _detect_device(self):
        """Detect best available device."""
        try:
            import torch
            if torch.cuda.is_available():
                return "cuda"
            elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
                return "mps"
        except ImportError:
            pass
        return "cpu"

    async def process_file_by_type(self, file_path: str, file_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process file based on its type."""
        file_extension = file_metadata.get('file_extension', '').lower()
        
        if file_extension == '.pdf':
            return self.process_pdf_mineru(file_path, file_metadata)
        elif file_extension in ['.mp3', '.wav', '.m4a']:
            return await self.process_audio_immediate(file_path, file_metadata)
        elif file_extension in [".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".3gp", ".ogv", ".m4v"]:
            return await self.process_video_immediate(file_path, file_metadata)
        elif file_extension == ".md":
            return self.process_markdown_direct(file_path, file_metadata)
        elif file_extension == ".txt":
            return self.process_text_immediate(file_path, file_metadata)
        elif file_extension in [".ppt", ".pptx"]:
            return self.process_presentation_immediate(file_path, file_metadata)
        else:
            return {
                'content': f"File type {file_extension} is supported but no immediate processing available.",
                'metadata': {},
                'features_available': [],
                'processing_time': 'immediate'
            }

    def process_pdf_mineru(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """PDF text extraction using MinerU API."""
        try:
            self.logger.info(f"Starting MinerU processing of {file_path}")
            start_time = time.time()

            # Check if MinerU API is available
            if not self.check_mineru_health():
                self.logger.warning("MinerU API not available, falling back to PyMuPDF")
                return self.process_pdf_pymupdf_fallback(file_path, file_metadata)

            # Generate clean filename for PDF
            original_filename = file_metadata.get('filename', 'document')
            base_title = original_filename.rsplit('.', 1)[0] if '.' in original_filename else original_filename
            
            # Import clean_title function
            try:
                from ..utils.helpers import clean_title
                clean_pdf_title = clean_title(base_title)
            except ImportError:
                clean_pdf_title = base_title

            # Call MinerU API
            mineru_result = self._call_mineru_api(file_path)
            
            # Extract the results for the first (and likely only) document
            results = mineru_result.get('results', {})
            if not results:
                raise Exception("No results returned from MinerU API")
            
            # Get the first document result (assuming single PDF)
            doc_key = next(iter(results.keys()))
            doc_result = results[doc_key]
            
            # Extract markdown content
            md_content = doc_result.get('md_content', '')
            
            # Extract images if available
            images = doc_result.get('images', {})
            
            # Save content to temporary file for processing
            temp_dir = tempfile.mkdtemp(suffix='_mineru_output')
            
            # Save markdown content
            md_file_path = os.path.join(temp_dir, f"{doc_key}.md")
            with open(md_file_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            # Save images if any
            image_files = []
            for img_name, img_data in images.items():
                if img_data.startswith('data:image/'):
                    # Handle base64 encoded images
                    header, data = img_data.split(',', 1)
                    img_bytes = base64.b64decode(data)
                    
                    img_path = os.path.join(temp_dir, img_name)
                    with open(img_path, 'wb') as f:
                        f.write(img_bytes)
                    image_files.append(img_path)
            
            # Get basic PDF metadata using PyMuPDF if available
            try:
                if fitz:
                    doc = fitz.open(file_path)
                    pdf_metadata = {
                        'page_count': doc.page_count,
                        'title': doc.metadata.get('title', ''),
                        'author': doc.metadata.get('author', ''),
                        'creation_date': doc.metadata.get('creationDate', ''),
                        'modification_date': doc.metadata.get('modDate', ''),
                        'processing_method': 'mineru_api',
                        'api_version': mineru_result.get('version', 'unknown'),
                        'backend': mineru_result.get('backend', 'pipeline'),
                        'has_markdown_content': bool(md_content),
                        'extracted_images_count': len(images),
                        'temp_output_dir': temp_dir
                    }
                    doc.close()
                else:
                    pdf_metadata = {
                        'processing_method': 'mineru_api',
                        'api_version': mineru_result.get('version', 'unknown'),
                        'backend': mineru_result.get('backend', 'pipeline'),
                        'has_markdown_content': bool(md_content),
                        'extracted_images_count': len(images),
                        'temp_output_dir': temp_dir
                    }
            except Exception as e:
                self.logger.warning(f"Could not extract PDF metadata: {e}")
                pdf_metadata = {
                    'processing_method': 'mineru_api',
                    'api_version': mineru_result.get('version', 'unknown'),
                    'metadata_error': str(e),
                    'has_markdown_content': bool(md_content),
                    'extracted_images_count': len(images),
                    'temp_output_dir': temp_dir
                }

            end_time = time.time()
            duration = end_time - start_time
            self.logger.info(f"MinerU processing completed in {duration:.2f} seconds")

            # Calculate features available
            features_available = ['advanced_pdf_extraction', 'markdown_conversion']
            if images:
                features_available.append('image_extraction')
            if 'table' in md_content.lower() or '|' in md_content:
                features_available.append('table_extraction')
            if any(marker in md_content for marker in ['$$', '$', '\\(']):
                features_available.append('formula_extraction')
            features_available.append('layout_analysis')

            result = {
                'content': md_content,  # Return markdown content directly
                'content_filename': f"{clean_pdf_title}.md",
                'metadata': pdf_metadata,
                'features_available': features_available,
                'processing_time': f'{duration:.2f}s',
                'mineru_extraction_result': {
                    'success': True,
                    'temp_output_dir': temp_dir,
                    'clean_title': clean_pdf_title,
                    'markdown_file': md_file_path,
                    'image_files': image_files,
                    'api_response': mineru_result
                }
            }

            return result

        except Exception as e:
            self.logger.warning(f"MinerU processing failed: {e}")
            # Fallback to PyMuPDF if MinerU fails
            return self.process_pdf_pymupdf_fallback(file_path, file_metadata)

    def process_pdf_pymupdf_fallback(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Fallback PDF text extraction using PyMuPDF when marker is not available."""
        try:
            self.logger.info(f"Using PyMuPDF fallback for {file_path}")

            if not fitz:
                raise Exception("PyMuPDF (fitz) is not available")

            doc = fitz.open(file_path)
            content = ""

            # Extract basic metadata
            pdf_metadata = {
                'page_count': doc.page_count,
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', ''),
                'processing_method': 'pymupdf_fallback'
            }

            # Extract text from all pages
            for page_num in range(doc.page_count):
                page = doc[page_num]
                content += f"\n=== Page {page_num + 1} ===\n"
                page_text = page.get_text()
                content += page_text

            doc.close()

            # Check if content extraction was successful
            if not content.strip():
                content = f"PDF document '{file_metadata['filename']}' appears to be image-based or empty. Text extraction may require OCR processing."

            return {
                "content": content,
                "metadata": pdf_metadata,
                "features_available": [
                    "advanced_pdf_extraction",
                    "figure_extraction",
                    "table_extraction",
                ],
                "processing_time": "immediate",
            }

        except Exception as e:
            raise Exception(f"PDF processing failed: {str(e)}")

    def _call_mineru_api(
        self,
        file_path: str,
        output_dir: str = "./output",
        lang_list: list = None,
        backend: str = "pipeline",
        parse_method: str = "auto",
        formula_enable: bool = True,
        table_enable: bool = True,
        server_url: str = "10.218.163.144",
        return_md: bool = True,
        return_middle_json: bool = False,
        return_model_output: bool = False,
        return_content_list: bool = False,
        return_images: bool = True,
        response_format_zip: bool = False,
        start_page_id: int = 0,
        end_page_id: int = 99999
    ) -> Dict[str, Any]:
        """Call the MinerU API to parse a PDF file."""
        if lang_list is None:
            lang_list = ['ch']
            
        # Validate file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")
        
        try:
            # Prepare the multipart form data
            files = {
                'files': (os.path.basename(file_path), open(file_path, 'rb'), 'application/pdf')
            }
            
            data = {
                'output_dir': output_dir,
                'lang_list': lang_list,
                'backend': backend,
                'parse_method': parse_method,
                'formula_enable': formula_enable,
                'table_enable': table_enable,
                'server_url': server_url,
                'return_md': return_md,
                'return_middle_json': return_middle_json,
                'return_model_output': return_model_output,
                'return_content_list': return_content_list,
                'return_images': return_images,
                'response_format_zip': response_format_zip,
                'start_page_id': start_page_id,
                'end_page_id': end_page_id
            }
            
            # Make the API request
            response = requests.post(
                self.mineru_parse_endpoint,
                files=files,
                data=data,
                timeout=300  # 5 minute timeout for large PDFs
            )
            
            # Close the file handle
            files['files'][1].close()
            
            response.raise_for_status()
            
            # Parse the JSON response
            result = response.json()
            
            return result
            
        except requests.exceptions.RequestException as e:
            self.logger.error(f"MinerU API request failed: {e}")
            raise Exception(f"MinerU API request failed: {e}")
        except json.JSONDecodeError as e:
            self.logger.error(f"Failed to parse MinerU API response: {e}")
            raise Exception(f"Invalid JSON response from MinerU API: {e}")
        except Exception as e:
            self.logger.error(f"MinerU API call failed: {e}")
            raise Exception(f"MinerU API call failed: {e}")

    async def process_audio_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Quick audio transcription using faster-whisper."""
        try:
            self.logger.info(f"Processing audio file: {file_path}")
            start_time = time.time()

            whisper_model = self.whisper_model
            if not whisper_model:
                return {
                    'content': f"Audio file '{file_metadata['filename']}' uploaded successfully. Transcription not available (faster-whisper not installed).",
                    'metadata': {'processing_method': 'audio_upload_only'},
                    'features_available': [],
                    'processing_time': 'immediate'
                }

            # Transcribe using faster-whisper
            segments, info = whisper_model.transcribe(file_path)
            
            transcript = ""
            for segment in segments:
                transcript += f"[{segment.start:.2f}s] {segment.text}\n"

            end_time = time.time()
            duration = end_time - start_time

            return {
                'content': transcript,
                'metadata': {
                    'duration': info.duration,
                    'language': info.language,
                    'language_probability': info.language_probability,
                    'processing_method': 'faster_whisper',
                    'model': 'base'
                },
                'features_available': ['audio_transcription', 'timestamped_segments'],
                'processing_time': f'{duration:.2f}s'
            }

        except Exception as e:
            raise Exception(f"Audio processing failed: {str(e)}")

    async def process_video_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process video by extracting audio and transcribing."""
        try:
            self.logger.info(f"Processing video file: {file_path}")
            start_time = time.time()

            # Extract audio from video using ffmpeg
            temp_audio_file = None
            try:
                import subprocess
                temp_audio_file = tempfile.mktemp(suffix='.wav')
                
                cmd = [
                    'ffmpeg', '-i', file_path, '-vn', '-acodec', 'pcm_s16le',
                    '-ar', '16000', '-ac', '1', temp_audio_file, '-y'
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    raise Exception(f"FFmpeg failed: {result.stderr}")

                # Transcribe the extracted audio
                audio_result = await self.process_audio_immediate(temp_audio_file, file_metadata)
                
                end_time = time.time()
                duration = end_time - start_time

                return {
                    'content': audio_result['content'],
                    'metadata': {
                        **audio_result['metadata'],
                        'original_type': 'video',
                        'video_processing_method': 'ffmpeg_audio_extraction'
                    },
                    'features_available': audio_result['features_available'] + ['video_to_audio_conversion'],
                    'processing_time': f'{duration:.2f}s'
                }

            finally:
                # Clean up temporary audio file
                if temp_audio_file and os.path.exists(temp_audio_file):
                    os.unlink(temp_audio_file)

        except Exception as e:
            raise Exception(f"Video processing failed: {str(e)}")

    def process_markdown_direct(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process markdown files directly."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            return {
                'content': content,
                'metadata': {
                    'processing_method': 'markdown_direct',
                    'character_count': len(content),
                    'line_count': len(content.splitlines())
                },
                'features_available': ['markdown_parsing', 'direct_content'],
                'processing_time': 'immediate'
            }

        except Exception as e:
            raise Exception(f"Markdown processing failed: {str(e)}")

    def process_text_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            return {
                'content': content,
                'metadata': {
                    'processing_method': 'text_direct',
                    'character_count': len(content),
                    'line_count': len(content.splitlines()),
                    'word_count': len(content.split())
                },
                'features_available': ['text_parsing', 'direct_content'],
                'processing_time': 'immediate'
            }

        except Exception as e:
            raise Exception(f"Text processing failed: {str(e)}")

    def process_presentation_immediate(self, file_path: str, file_metadata: Dict) -> Dict[str, Any]:
        """Process PowerPoint presentations."""
        try:
            content = ""
            
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                
                for i, slide in enumerate(prs.slides):
                    content += f"\n=== Slide {i + 1} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"

                return {
                    'content': content,
                    'metadata': {
                        'processing_method': 'python_pptx',
                        'slide_count': len(prs.slides)
                    },
                    'features_available': ['text_extraction', 'slide_structure'],
                    'processing_time': 'immediate'
                }

            except ImportError:
                # Fallback message if python-pptx is not available
                return {
                    'content': f"PowerPoint file '{file_metadata['filename']}' uploaded successfully. Text extraction not available (python-pptx not installed).",
                    'metadata': {'processing_method': 'presentation_upload_only'},
                    'features_available': [],
                    'processing_time': 'immediate'
                }

        except Exception as e:
            raise Exception(f"Presentation processing failed: {str(e)}") 